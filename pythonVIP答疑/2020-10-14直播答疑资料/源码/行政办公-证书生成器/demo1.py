#encoding=utf-8
#教育机构 ：马士兵教育
#讲    师：杨淑娟

from  pptx import Presentation
from openpyxl import  load_workbook
import  datetime
workbook=load_workbook('学员名单.xlsx')
sheet=workbook.active
rows=sheet.rows
lst=[]
for row in rows:
    stu_lst=[]
    for cell in row:
        stu_lst.append(cell.value)
    lst.append(stu_lst)
lst.remove(lst[0])
for item in lst:
    year_month=item[1].split('-')
    item.remove(item[1])
    item.insert(1,year_month[0])
    item.insert(2,year_month[1])

newdate=datetime.datetime.now()
for item in lst:
    item.append(newdate.month)
    item.append(newdate.day)
    item.append(newdate.year)
prs=Presentation('证书模板.pptx')
slide=prs.slides[0]
shapes=slide.shapes
for item in lst:
    shapes[0].text_frame.text=item[0]
    shapes[1].text_frame.text=str(item[1])
    shapes[2].text_frame.text=str(item[2])
    shapes[3].text_frame.text=item[3]
    shapes[4].text_frame.text=str(item[4])
    shapes[5].text_frame.text=str(item[5])
    shapes[6].text_frame.text=str(item[6])
    prs.save('student/'+str(item[0])+'.pptx')