from pptx import Presentation
from docx import Document

pre=Presentation('第7站-夫妻站.pptx')
lst=[]      # 用于存储每一个段落
for slide in pre.slides:    # 从幻灯片列表中遍历每一张幻灯片
    # 获取形状shape
    for shape in slide.shapes:
        # 判断形状是否存在文字
        if shape.has_text_frame:
            # 获取文字框
            text_frame=shape.text_frame
            for para in text_frame.paragraphs:
                lst.append(para.text)

# print(lst)
# 创建文档对象
doc=Document()
for item in lst:
    doc.add_paragraph(item)
doc.save('第7站-夫妻站.docx')